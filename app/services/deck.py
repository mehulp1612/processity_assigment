"""The analysis deck: a PPTX with real charts drawn from the books.

Every number and every chart comes from ``analytics.sales_report`` — the deck is
a *view* of the books, not a second calculation. If a figure on a slide is wrong,
it is wrong in the database.

Charts are rendered with matplotlib to PNG and embedded. That is deliberate: a
native PowerPoint chart object would carry its own copy of the data and could
drift from the invoice totals, and matplotlib gives control over how a zero-sale
day is drawn.
"""

from __future__ import annotations

from datetime import datetime
from typing import Optional

from . import analytics
from .documents import OUT_DIR
from .inventory import low_stock
from .memory import get_shop

# A restrained palette: one accent for money, one for tax, greys for structure.
_INK = "#1f2933"
_ACCENT = "#1f6f5c"
_ACCENT_2 = "#c98a2e"
_MUTED = "#8a94a0"
_SERIES = ["#1f6f5c", "#c98a2e", "#3f6d9e", "#8a5a83", "#6b7a52", "#a8553a"]


def _save(fig, name: str) -> str:
    import matplotlib.pyplot as plt

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f".chart-{name}.png"
    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return str(path)


def _style(ax) -> None:
    for side in ("top", "right"):
        ax.spines[side].set_visible(False)
    for side in ("left", "bottom"):
        ax.spines[side].set_color(_MUTED)
    ax.tick_params(colors=_INK, labelsize=8)
    ax.set_axisbelow(True)


def _chart_daily(report: dict) -> Optional[str]:
    import matplotlib.pyplot as plt

    series = report["by_day"]
    if not series:
        return None
    labels = [datetime.fromisoformat(d["day"]).strftime("%d %b") for d in series]
    values = [float(d["total"]) for d in series]

    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    bars = ax.bar(labels, values, color=_ACCENT, width=0.62)
    for bar, value in zip(bars, values):
        if value:                       # a "0" label over an empty bar is noise
            ax.annotate(f"{value:,.0f}", (bar.get_x() + bar.get_width() / 2, value),
                        ha="center", va="bottom", fontsize=7.5, color=_INK)
    ax.set_ylabel("Sales (Rs.)", fontsize=8, color=_INK)
    ax.grid(axis="y", color="#e6e9ed", linewidth=0.8)
    _style(ax)
    return _save(fig, "daily")


def _chart_mix(report: dict) -> Optional[str]:
    import matplotlib.pyplot as plt

    mix = [m for m in report["payment_mix"] if float(m["amount"]) > 0]
    if not mix:
        return None
    fig, ax = plt.subplots(figsize=(3.6, 3.2))
    ax.pie(
        [float(m["amount"]) for m in mix],
        labels=[str(m["mode"]).upper() for m in mix],
        autopct="%1.0f%%", startangle=90, counterclock=False,
        colors=_SERIES[: len(mix)],
        textprops={"fontsize": 8, "color": _INK},
        wedgeprops={"edgecolor": "white", "linewidth": 1.5},
    )
    ax.axis("equal")
    return _save(fig, "mix")


def _chart_top(report: dict) -> Optional[str]:
    import matplotlib.pyplot as plt

    items = report["top_items"][:6]
    if not items:
        return None
    names = [(i["name"] if len(i["name"]) <= 26 else i["name"][:25] + "…") for i in items][::-1]
    revenue = [float(i["revenue"]) for i in items][::-1]

    fig, ax = plt.subplots(figsize=(7.4, 3.2))
    ax.barh(names, revenue, color=_ACCENT_2, height=0.6)
    for y, value in enumerate(revenue):
        ax.annotate(f"{value:,.0f}", (value, y), xytext=(4, 0), textcoords="offset points",
                    va="center", fontsize=7.5, color=_INK)
    ax.set_xlabel("Revenue ex-GST (Rs.)", fontsize=8, color=_INK)
    ax.grid(axis="x", color="#e6e9ed", linewidth=0.8)
    _style(ax)
    return _save(fig, "top")


def build_analysis_pptx(start: Optional[str] = None, end: Optional[str] = None) -> dict:
    """Build the analysis deck for a shop-local date range. Returns the file path."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    report = analytics.sales_report(start, end, top=10)
    shop = get_shop() or {}
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / f"sales-{report['start']}-to-{report['end']}.pptx"

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    blank = prs.slide_layouts[6]
    ink = RGBColor.from_string(_INK.lstrip("#"))
    accent = RGBColor.from_string(_ACCENT.lstrip("#"))
    muted = RGBColor.from_string(_MUTED.lstrip("#"))

    def text(slide, body, x, y, w, h, size, *, bold=False, color=ink):
        frame = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
        frame.word_wrap = True
        run = frame.paragraphs[0].add_run()
        run.text = body
        run.font.size, run.font.bold, run.font.color.rgb = Pt(size), bold, color

    def pretty(value: str) -> str:
        return datetime.fromisoformat(value).strftime("%d %b %Y")

    period = f"{pretty(report['start'])} - {pretty(report['end'])}"
    slides = 0

    # 1. Title
    slide = prs.slides.add_slide(blank); slides += 1
    text(slide, shop.get("name", "Kirana Store"), 0.9, 2.4, 11.5, 1.0, 40, bold=True)
    text(slide, "Sales & GST Analysis", 0.9, 3.4, 11.5, 0.7, 22, color=accent)
    text(slide, period, 0.9, 4.15, 11.5, 0.5, 14, color=muted)
    if shop.get("gstin"):
        text(slide, f"GSTIN {shop['gstin']}", 0.9, 6.4, 11.5, 0.4, 10, color=muted)

    # 2. Headline numbers
    slide = prs.slides.add_slide(blank); slides += 1
    text(slide, "At a glance", 0.7, 0.5, 8, 0.6, 26, bold=True)
    text(slide, period, 0.7, 1.05, 8, 0.4, 11, color=muted)
    for i, (label, value) in enumerate([
        ("Total sales", f"Rs. {report['total']:,.0f}"),
        ("Bills", f"{report['bills']}"),
        ("GST collected", f"Rs. {report['gst']:,.2f}"),
        ("Gross margin", f"Rs. {report['gross_margin']:,.0f}"),
        ("Margin", f"{report['margin_pct']:.1f}%"),
        ("Credit outstanding", f"Rs. {report['khata']['outstanding_total']:,.0f}"),
    ]):
        x, y = 0.7 + (i % 3) * 4.15, 2.0 + (i // 3) * 2.1
        text(slide, value, x, y, 3.9, 0.9, 30, bold=True, color=accent)
        text(slide, label, x, y + 0.85, 3.9, 0.4, 12, color=muted)

    # 3-5. Charts
    for title, subtitle, maker in (
        ("Sales by day", "Daily takings including GST", _chart_daily),
        ("How customers paid", "Share of total billed value", _chart_mix),
        ("Best sellers", "By revenue excluding GST", _chart_top),
    ):
        png = maker(report)
        if png is None:
            continue
        slide = prs.slides.add_slide(blank); slides += 1
        text(slide, title, 0.7, 0.45, 9, 0.6, 24, bold=True)
        text(slide, subtitle, 0.7, 1.0, 9, 0.4, 11, color=muted)
        slide.shapes.add_picture(png, Inches(0.9), Inches(1.75), height=Inches(4.6))

    # 6. GST by slab
    slabs = report["slabs"] or []
    slide = prs.slides.add_slide(blank); slides += 1
    text(slide, "GST by slab", 0.7, 0.45, 9, 0.6, 24, bold=True)
    text(slide, "Intra-state supply: CGST and SGST at half the rate each",
         0.7, 1.0, 9, 0.4, 11, color=muted)
    table = slide.shapes.add_table(
        len(slabs) + 2, 4, Inches(0.9), Inches(1.8), Inches(9.0),
        Inches(0.42 * (len(slabs) + 2)),
    ).table
    for col, head in enumerate(["Slab", "Taxable value", "CGST", "SGST"]):
        table.cell(0, col).text = head
    for row, data in enumerate(slabs, start=1):
        for col, value in enumerate([
            f"{data['gst_rate']:g}%", f"Rs. {data['taxable']:,.2f}",
            f"Rs. {data['cgst']:,.2f}", f"Rs. {data['sgst']:,.2f}",
        ]):
            table.cell(row, col).text = value
    for col, value in enumerate([
        "Total", f"Rs. {report['subtotal']:,.2f}",
        f"Rs. {report['cgst']:,.2f}", f"Rs. {report['sgst']:,.2f}",
    ]):
        cell = table.cell(len(slabs) + 1, col)
        cell.text = value
        cell.text_frame.paragraphs[0].runs[0].font.bold = True
    for row in range(len(slabs) + 2):
        for col in range(4):
            for para in table.cell(row, col).text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(12)

    # 7. What to act on
    slide = prs.slides.add_slide(blank); slides += 1
    text(slide, "What to act on", 0.7, 0.45, 9, 0.6, 24, bold=True)

    bullets: list[str] = []
    reorder = low_stock()
    if reorder:
        bullets.append(f"Reorder now - {len(reorder)} item(s) at or below reorder level:")
        bullets += [f"    - {p['name']}: {p['qty']:g} {p['unit']} left "
                    f"(reorder at {p['reorder_level']:g})" for p in reorder[:6]]
    else:
        bullets.append("Stock is above reorder level on every item.")

    outstanding = report["khata"]["outstanding_total"]
    bullets.append(
        f"Credit outstanding: Rs. {outstanding:,.0f} across the khata book."
        if outstanding else "No credit outstanding."
    )
    if report["top_items"]:
        best = report["top_items"][0]
        bullets.append(f"Best seller: {best['name']} - Rs. {best['revenue']:,.0f} revenue, "
                       f"Rs. {best['margin']:,.0f} margin.")
    bullets.append(f"Overall margin {report['margin_pct']:.1f}% on Rs. "
                   f"{report['revenue_ex_gst']:,.0f} of sales (ex-GST).")

    frame = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.0)).text_frame
    frame.word_wrap = True
    for i, line in enumerate(bullets):
        para = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = line
        indented = line.startswith("    ")
        run.font.size = Pt(13 if indented else 15)
        run.font.color.rgb = muted if indented else ink
        para.space_after = Pt(8)

    prs.save(str(path))

    for stale in OUT_DIR.glob(".chart-*.png"):   # PNGs were only an embedding step
        stale.unlink(missing_ok=True)

    return {
        "ok": True,
        "path": str(path),
        "start": report["start"],
        "end": report["end"],
        "slides": slides,
        "bills": report["bills"],
        "total": report["total"],
    }
